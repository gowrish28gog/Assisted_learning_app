from pypdf import PdfReader
import convertapi
import streamlit as st
from io import BytesIO
import docx
from pptx import Presentation

# Set up ConvertAPI credentials
convertapi.api_credentials =''


# Function to extract text from PDF

def extract_text_from_pdf(pdf_file: BytesIO) -> str:
    """
    Extracts plain text from a PDF file.

    :param pdf_file: A file-like object for the PDF file.
    :return: Extracted text as a string.
    """
    try:
        reader = PdfReader(pdf_file)
        text = "\n".join([page.extract_text() for page in reader.pages if page.extract_text()])
        return text.strip()
    except Exception as e:
        st.error(f"Error extracting text from PDF: {e}")
        return ""


# Function to convert docx or pptx to PDF

def convert_and_extract_pdf(file: BytesIO, file_extension: str) -> str:
    """
    Converts DOCX or PPTX file to PDF and extracts the text from it.

    :param file: A file-like object for the input DOCX/PPTX file.
    :param file_extension: The extension of the uploaded file (either 'docx' or 'pptx').
    :return: Extracted text as a string.
    """
    # Convert DOCX/PPTX to PDF using ConvertAPI
    result = convertapi.convert('pdf', {'File': file}, from_format=file_extension)
    
    # Save the converted PDF
    pdf_path = result.save_files('temp')[0]

    # Extract text from the converted PDF
    with open(pdf_path, 'rb') as f:
        pdf_file = BytesIO(f.read())
    return extract_text_from_pdf(pdf_file)


# Function to extract text from docx

def extract_text_from_docx(docx_file: BytesIO) -> str:
    """
    Extracts plain text from a DOCX file.

    :param docx_file: A file-like object for the DOCX file.
    :return: Extracted text as a string.
    """
    try:
        doc = docx.Document(docx_file)
        text = "\n".join([paragraph.text for paragraph in doc.paragraphs])
        return text.strip()
    except Exception as e:
        st.error(f"Error extracting text from DOCX: {e}")
        return ""


# Function to extract text from pptx

def extract_text_from_pptx(pptx_file: BytesIO) -> str:
    """
    Extracts plain text from a PPTX file.

    :param pptx_file: A file-like object for the PPTX file.
    :return: Extracted text as a string.
    """
    try:
        presentation = Presentation(pptx_file)
        text = []
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return "\n".join(text).strip()
    except Exception as e:
        st.error(f"Error extracting text from PPTX: {e}")
        return ""


# Function to extract text according to document type

def extract_text(uploaded_file) -> str:
    """
    Extracts text from a PDF, DOCX, or PPTX file. If the file is not a PDF, it converts it to PDF first.

    :param uploaded_file: The uploaded file object (either PDF, DOCX, or PPTX).
    :return: Extracted text as a string.
    """
    # Get the file name and extension from the uploaded file
    file_name = uploaded_file.name
    file_extension = file_name.split('.')[-1].lower()

    # Read the file as a BytesIO object
    file_content = BytesIO(uploaded_file.read())

    # Process based on the file type
    if file_extension == 'pdf':
        # If the file is already a PDF, directly extract text
        return extract_text_from_pdf(file_content)
    elif file_extension == 'docx':
        # If the file is a DOCX, extract text
        return extract_text_from_docx(file_content)
    elif file_extension == 'pptx':
        # If the file is a PPTX, extract text
        return extract_text_from_pptx(file_content)
    else:
        raise ValueError("Unsupported file type. Please upload a PDF, DOCX, or PPTX file.")